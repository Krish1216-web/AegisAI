import os
from pptx import Presentation
from loguru import logger

from app.services.extractors.base import BaseDocumentExtractor, ExtractedDocument, PageData
from app.core.document_exceptions import InvalidFile

class PPTXExtractor(BaseDocumentExtractor):
    def extract(self, file_path: str) -> ExtractedDocument:
        if not os.path.exists(file_path):
            raise InvalidFile("PPTX file does not exist on disk.")

        try:
            prs = Presentation(file_path)
            pages_data = []
            full_text_parts = []
            
            for idx, slide in enumerate(prs.slides):
                slide_num = idx + 1
                slide_text_parts = []
                
                # Try getting the title shape specifically if present
                title = ""
                try:
                    if slide.shapes.title:
                        title = slide.shapes.title.text.strip()
                        if title:
                            slide_text_parts.append(f"[Title] {title}")
                except Exception:
                    pass

                # Iterate other shapes to extract text
                for shape in slide.shapes:
                    # Skip if shape is the title shape (we already parsed it)
                    try:
                        if slide.shapes.title and shape == slide.shapes.title:
                            continue
                    except Exception:
                        pass
                        
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())
                        
                slide_text = "\n".join(slide_text_parts)
                pages_data.append(
                    PageData(
                        page_number=slide_num,
                        text=slide_text,
                        metadata={"slide_index": idx}
                    )
                )
                full_text_parts.append(f"--- Slide {slide_num} ---\n{slide_text}")

            full_text = "\n\n".join(full_text_parts).strip()
            char_count = len(full_text)
            word_count = len(full_text.split())

            return ExtractedDocument(
                text=full_text,
                pages=pages_data,
                sections=[],
                metadata={},
                page_count=len(prs.slides),
                character_count=char_count,
                word_count=word_count
            )

        except Exception as e:
            logger.error(f"Failed to parse PPTX: {e}")
            raise InvalidFile(f"Failed to parse PPTX document: {str(e)}")
