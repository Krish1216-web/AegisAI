import pytest
import os
import tempfile
from pptx import Presentation

from app.services.extractors.pptx import PPTXExtractor

def test_pptx_extractor_success():
    extractor = PPTXExtractor()

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as f:
        temp_path = f.name
        
    try:
        prs = Presentation()
        # Add slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title and text
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        subtitle = slide.placeholders[1]
        subtitle.text = "Q3 Performance Review"
        
        prs.save(temp_path)
        
        res = extractor.extract(temp_path)
        
        assert res.page_count == 1
        assert "Executive Summary" in res.text
        assert "Q3 Performance Review" in res.text
        assert len(res.pages) == 1
        assert res.pages[0].page_number == 1
        
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)
